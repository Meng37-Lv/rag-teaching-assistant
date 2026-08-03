from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PPT = PROJECT_ROOT / "ppt" / "course.pptx"
DEFAULT_PPT_DIR = PROJECT_ROOT / "ppt"
DEFAULT_OUTPUT = PROJECT_ROOT / "data" / "text.txt"
CHAPTER_ORDER = {
    "一": 1,
    "二": 2,
    "三": 3,
    "四": 4,
    "五": 5,
    "六": 6,
    "七": 7,
    "八": 8,
    "九": 9,
    "十": 10,
}


def iter_shape_text(shape) -> list[str]:
    texts: list[str] = []

    if hasattr(shape, "text") and shape.text:
        text = shape.text.strip()
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                texts.append(" | ".join(cells))

    if getattr(shape, "has_group_shape", False):
        for child in shape.shapes:
            texts.extend(iter_shape_text(child))

    return texts


def extract_slide_text(slide) -> str:
    slide_texts: list[str] = []
    for shape in slide.shapes:
        slide_texts.extend(iter_shape_text(shape))
    return "\n".join(slide_texts).strip()


def extract_ppt_text(ppt_paths: list[Path]) -> str:
    if not ppt_paths:
        raise FileNotFoundError("No PPT files found.")

    parts: list[str] = []
    page_number = 1

    for ppt_path in ppt_paths:
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")

        presentation = Presentation(str(ppt_path))
        for slide in presentation.slides:
            slide_text = extract_slide_text(slide)
            parts.append(f"===== 第{page_number}页 =====\n\n{slide_text}")
            page_number += 1

    return "\n\n\n".join(parts).strip() + "\n"


def resolve_ppt_paths(ppt_paths: list[Path] | None, ppt_dir: Path | None) -> list[Path]:
    if ppt_paths:
        return ppt_paths

    if DEFAULT_PPT.exists():
        return [DEFAULT_PPT]

    search_dir = ppt_dir or DEFAULT_PPT_DIR
    if search_dir.exists():
        return sorted(search_dir.glob("*.pptx"), key=ppt_sort_key)

    return [DEFAULT_PPT]


def ppt_sort_key(path: Path) -> tuple[int, str]:
    name = path.stem
    if name.startswith("第") and "章" in name:
        chapter_text = name[1 : name.index("章")]
        chapter_number = CHAPTER_ORDER.get(chapter_text)
        if chapter_number is not None:
            return chapter_number, name

    return 999, name


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract readable text from a course PPTX file.")
    parser.add_argument("--ppt", type=Path, nargs="*", help="Input PPTX path. Multiple files are supported.")
    parser.add_argument("--ppt-dir", type=Path, help="Directory containing PPTX files.")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT, help="Output text path.")
    args = parser.parse_args()

    ppt_paths = resolve_ppt_paths(args.ppt, args.ppt_dir)
    text = extract_ppt_text(ppt_paths)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(text, encoding="utf-8")
    print("PPT文本提取完成")


if __name__ == "__main__":
    main()
