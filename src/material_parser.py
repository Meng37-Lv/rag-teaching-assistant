from __future__ import annotations

from dataclasses import asdict, dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_MATERIAL_TYPES = {
    ".pptx": "PPTX",
    ".docx": "DOCX",
    ".md": "Markdown",
    ".txt": "TXT",
    ".pdf": "PDF",
}


class MaterialParseError(ValueError):
    """可安全展示给用户的材料解析错误。"""


@dataclass(frozen=True)
class ParsedMaterial:
    material_name: str
    material_type: str
    extracted_text: str

    def to_dict(self) -> dict[str, str]:
        return asdict(self)


def parse_material_file(file_path: str | Path) -> ParsedMaterial:
    """解析受支持的本地材料文件，不写入项目数据目录。"""

    path = Path(file_path)
    material_type = SUPPORTED_MATERIAL_TYPES.get(path.suffix.lower())
    if material_type is None:
        raise MaterialParseError("不支持该材料格式，请使用 PPTX、DOCX、Markdown 或 TXT 文件。")
    if not path.is_file():
        raise MaterialParseError("材料文件不存在或无法读取。")

    try:
        if path.suffix.lower() == ".pptx":
            extracted_text = _extract_pptx(path)
        elif path.suffix.lower() == ".docx":
            extracted_text = _extract_docx(path)
        elif path.suffix.lower() == ".pdf":
            extracted_text = "\n".join((page.extract_text() or "").strip() for page in PdfReader(str(path)).pages).strip()
        else:
            extracted_text = _extract_text_file(path)
    except MaterialParseError:
        raise
    except Exception:
        raise MaterialParseError("材料解析失败，请确认文件未损坏且格式正确。") from None

    return ParsedMaterial(
        material_name=path.name,
        material_type=material_type,
        extracted_text=extracted_text,
    )


def parse_text_material(
    text: str,
    material_name: str = "纯文本材料",
) -> ParsedMaterial:
    """将调用方提供的纯文本转换为统一材料结果。"""

    if not isinstance(text, str):
        raise MaterialParseError("纯文本材料必须是字符串。")
    safe_name = material_name.strip() if isinstance(material_name, str) else ""
    if not safe_name:
        safe_name = "纯文本材料"
    return ParsedMaterial(
        material_name=safe_name,
        material_type="纯文本",
        extracted_text=text,
    )


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    pages: list[str] = []
    for page_number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            shape_text = getattr(shape, "text", "")
            if shape_text and shape_text.strip():
                texts.append(shape_text.strip())
        page_text = "\n".join(texts)
        pages.append(f"===== 第{page_number}页 =====\n{page_text}".rstrip())
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    document = Document(str(path))
    blocks = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                blocks.append("\t".join(cells))
    return "\n".join(blocks)


def _extract_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise MaterialParseError("文本文件编码无法识别，请将文件保存为 UTF-8。")
