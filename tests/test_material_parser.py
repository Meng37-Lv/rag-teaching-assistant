from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from docx import Document
from pptx import Presentation

from src.material_parser import (
    MaterialParseError,
    parse_material_file,
    parse_text_material,
)


class MaterialParserTests(unittest.TestCase):
    def test_parse_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "最小汇报.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "汇报标题"
            slide.placeholders[1].text = "PPT正文"
            presentation.save(path)

            result = parse_material_file(path)

        self.assertEqual(result.material_name, "最小汇报.pptx")
        self.assertEqual(result.material_type, "PPTX")
        self.assertIn("===== 第1页 =====", result.extracted_text)
        self.assertIn("汇报标题", result.extracted_text)
        self.assertIn("PPT正文", result.extracted_text)

    def test_parse_docx(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "最小汇报.docx"
            document = Document()
            document.add_paragraph("DOCX正文")
            document.save(path)

            result = parse_material_file(path)

        self.assertEqual(result.material_type, "DOCX")
        self.assertEqual(result.extracted_text, "DOCX正文")

    def test_parse_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "最小汇报.md"
            path.write_text("# 标题\n\nMarkdown正文", encoding="utf-8")
            result = parse_material_file(path)

        self.assertEqual(result.material_type, "Markdown")
        self.assertEqual(result.extracted_text, "# 标题\n\nMarkdown正文")

    def test_parse_txt(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "最小汇报.txt"
            path.write_text("TXT正文", encoding="utf-8")
            result = parse_material_file(path)

        self.assertEqual(result.material_type, "TXT")
        self.assertEqual(result.extracted_text, "TXT正文")

    def test_parse_plain_text(self) -> None:
        result = parse_text_material("纯文本正文", "课堂汇报")

        self.assertEqual(
            result.to_dict(),
            {
                "material_name": "课堂汇报",
                "material_type": "纯文本",
                "extracted_text": "纯文本正文",
            },
        )

    def test_failure_message_does_not_expose_path(self) -> None:
        private_path = Path("C:/private/student/不存在.docx")

        with self.assertRaises(MaterialParseError) as context:
            parse_material_file(private_path)

        self.assertEqual(str(context.exception), "材料文件不存在或无法读取。")
        self.assertNotIn(str(private_path), str(context.exception))

    def test_corrupted_docx_returns_safe_error(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "学生汇报.docx"
            path.write_text("这不是有效的DOCX文件", encoding="utf-8")

            with self.assertRaises(MaterialParseError) as context:
                parse_material_file(path)

        self.assertEqual(
            str(context.exception),
            "材料解析失败，请确认文件未损坏且格式正确。",
        )
        self.assertNotIn(str(path), str(context.exception))

    def test_corrupted_pdf_returns_safe_error(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "empty.pdf"
            path.write_bytes(b"not-a-pdf")
            with self.assertRaises(MaterialParseError) as context:
                parse_material_file(path)
        self.assertEqual(str(context.exception), "材料解析失败，请确认文件未损坏且格式正确。")


if __name__ == "__main__":
    unittest.main()
